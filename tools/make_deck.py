from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = Path(__file__).resolve().parent.parent / "CarePulse_Pitch_Deck.pptx"
FONT = "Segoe UI"
MONO = "Consolas"


def C(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


BG = C("0B1120")
PANEL = C("121C30")
PANEL2 = C("0E1729")
ROW_ALT = C("10192C")
ACCENT = C("38BDF8")
BLUE = C("2563EB")
WHITE = C("F8FAFC")
TEXT = C("DBE4F0")
MUTED = C("93A4BC")
FAINT = C("5B6B85")
ROSE = C("FB7185")
AMBER = C("FBBF24")
GREEN = C("34D399")
LINE = C("233350")

TIER_COLORS = {"rose": ROSE, "amber": AMBER, "green": GREEN, "sky": ACCENT}

SW = Inches(13.333)
SH = Inches(7.5)
MARGIN = Inches(0.55)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def run(p, text, size=16, color=TEXT, bold=False, name=FONT, italic=False):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = name
    f.color.rgb = color
    return r


def rect(slide, x, y, w, h, fill, line_color=None, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(kind, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def oval_outline(slide, x, y, w, h, color, width_pt):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    o.fill.background()
    o.line.color.rgb = color
    o.line.width = Pt(width_pt)
    o.shadow.inherit = False
    return o


def footer(slide, idx):
    ln = rect(slide, MARGIN, SH - Inches(0.52), SW - MARGIN * 2, Pt(1), LINE)
    ln.shadow.inherit = False
    tf = box(slide, MARGIN, SH - Inches(0.46), Inches(6), Inches(0.3))
    p = tf.paragraphs[0]
    run(p, "P E R S O N N E L A I", 10, FAINT, True)
    tf2 = box(slide, SW - MARGIN - Inches(1.2), SH - Inches(0.46), Inches(1.2), Inches(0.3))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run(p2, "%02d" % idx, 10, FAINT, True)


def header(slide, kicker, title, idx, accent=ACCENT):
    tf = box(slide, MARGIN, Inches(0.42), SW - MARGIN * 2, Inches(0.32))
    p = tf.paragraphs[0]
    run(p, kicker.upper(), 12, accent, True)
    tf = box(slide, MARGIN, Inches(0.74), SW - MARGIN * 2, Inches(0.72))
    p = tf.paragraphs[0]
    run(p, title, 31, WHITE, True)
    rect(slide, MARGIN, Inches(1.44), Inches(0.9), Pt(3), accent)
    footer(slide, idx)


def bullets(tf, items, size=15.5, gap=8, glyph="▸", glyph_color=ACCENT):
    first = True
    for it in items:
        p = para(tf, first)
        first = False
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        run(p, glyph + "  ", size, glyph_color, True)
        if isinstance(it, tuple):
            lead, rest = it
            run(p, lead + " ", size, WHITE, True)
            run(p, rest, size, TEXT, False)
        else:
            run(p, it, size, TEXT, False)


def chips_row(slide, items, y, colors=None, centered=False):
    widths = []
    for t in items:
        widths.append(Inches(0.42 + len(t) * 0.082))
    gap = Inches(0.18)
    total = sum(widths, Inches(0)) + gap * (len(items) - 1)
    x = int((SW - total) / 2) if centered else MARGIN
    for i, t in enumerate(items):
        colr = colors[i % len(colors)] if colors else ACCENT
        c = rect(slide, x, y, widths[i], Inches(0.42), PANEL, LINE, radius=0.5)
        tfc = c.text_frame
        tfc.word_wrap = False
        tfc.margin_left = Inches(0.06)
        tfc.margin_right = Inches(0.06)
        tfc.margin_top = 0
        tfc.margin_bottom = 0
        tfc.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tfc.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, t, 12, colr, True)
        x = x + widths[i] + gap


def callout(slide, text, y, color=AMBER, h=Inches(0.62)):
    pan = rect(slide, MARGIN, y, SW - MARGIN * 2, h, PANEL2, radius=0.18)
    pan.shadow.inherit = False
    rect(slide, MARGIN, y, Pt(3.2), h, color)
    tf = pan.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run(p, text, 14, color, True)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def slide_content(idx, kicker, title, items, note, glyph="▸", glyph_color=ACCENT,
                  callout_text=None, callout_color=AMBER, chips=None, chip_colors=None,
                  body_size=15.5):
    s = new_slide()
    header(s, kicker, title, idx)
    tf = box(s, MARGIN, Inches(1.82), SW - MARGIN * 2, Inches(3.55))
    bullets(tf, items, size=body_size, glyph=glyph, glyph_color=glyph_color)
    y_call = SH - Inches(1.78)
    if callout_text:
        callout(s, callout_text, y_call, callout_color)
        y_chips = SH - Inches(0.95)
    else:
        y_chips = SH - Inches(1.05)
    if chips:
        chips_row(s, chips, y_chips, chip_colors)
    notes(s, note)
    return s


def slide_two_col(idx, kicker, title, left_head, left_items, right_head, right_items, note):
    s = new_slide()
    header(s, kicker, title, idx)
    col_w = int((SW - MARGIN * 2 - Inches(0.4)) / 2)
    for i, (h_txt, items) in enumerate([(left_head, left_items), (right_head, right_items)]):
        x = MARGIN + i * (col_w + Inches(0.4))
        pan = rect(s, x, Inches(1.85), col_w, Inches(4.35), PANEL2, radius=0.09)
        pan.shadow.inherit = False
        rect(s, x, Inches(1.85), col_w, Pt(3), ACCENT if i == 0 else BLUE)
        tf = box(s, x + Inches(0.28), Inches(2.08), col_w - Inches(0.56), Inches(0.4))
        p = tf.paragraphs[0]
        run(p, h_txt.upper(), 14, ACCENT if i == 0 else C("60A5FA"), True)
        tfb = box(s, x + Inches(0.28), Inches(2.58), col_w - Inches(0.56), Inches(3.4))
        bullets(tfb, items, size=13.5, gap=7)
    notes(s, note)
    return s


def slide_table(idx, kicker, title, headers, rows, col_widths, note, row_h=Inches(0.5),
                font_size=12.5, under_items=None):
    s = new_slide()
    header(s, kicker, title, idx)
    n_r = len(rows) + 1
    n_c = len(headers)
    total_w = sum(col_widths, Inches(0))
    tbl_h = Inches(0.46) + row_h * len(rows)
    gf = s.shapes.add_table(n_r, n_c, MARGIN, Inches(1.78), total_w, tbl_h)
    t = gf.table
    t.first_row = False
    t.horz_banding = False
    for i, w in enumerate(col_widths):
        t.columns[i].width = w
    t.rows[0].height = Inches(0.46)
    for i in range(len(rows)):
        t.rows[i + 1].height = row_h

    def style_cell(cell, val, is_header, r_i, col_i):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.02)
        cell.margin_bottom = Inches(0.02)
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = PANEL
        else:
            cell.fill.fore_color.rgb = PANEL2 if r_i % 2 == 1 else ROW_ALT
        tfc = cell.text_frame
        tfc.word_wrap = True
        p = tfc.paragraphs[0]
        if is_header:
            run(p, val.upper(), 11.5, ACCENT, True)
            return
        color = TEXT
        bold = False
        if isinstance(val, tuple):
            val, key = val
            color = TIER_COLORS.get(key, TEXT)
            bold = True
        elif val.startswith("HIGH"):
            color, bold = ROSE, True
        run(p, val, font_size, color, bold)

    for c, htxt in enumerate(headers):
        style_cell(t.cell(0, c), htxt, True, 0, c)
    for r_i, r in enumerate(rows, start=1):
        for c, val in enumerate(r):
            style_cell(t.cell(r_i, c), str(val), False, r_i, c)
    if under_items:
        tfu = box(s, MARGIN, Inches(1.78) + tbl_h + Inches(0.22), SW - MARGIN * 2, Inches(1.4))
        bullets(tfu, under_items, size=13, gap=5, glyph="•", glyph_color=C("60A5FA"))
    notes(s, note)
    return s


def slide_diagram(idx, kicker, title, mono_lines, points, note):
    s = new_slide()
    header(s, kicker, title, idx)
    pan = rect(s, MARGIN, Inches(1.85), SW - MARGIN * 2, Inches(2.6), PANEL2, radius=0.07)
    pan.shadow.inherit = False
    tfd = box(s, MARGIN + Inches(0.4), Inches(2.05), SW - MARGIN * 2 - Inches(0.8), Inches(2.25))
    first = True
    for ln in mono_lines:
        p = para(tfd, first)
        first = False
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)
        if ln.strip() in ("↓", "▼"):
            run(p, ln, 14, FAINT, True, MONO)
        else:
            run(p, ln, 15, TEXT, True, MONO)
    tfb = box(s, MARGIN, Inches(4.75), SW - MARGIN * 2, Inches(1.6))
    bullets(tfb, points, size=15, gap=8)
    notes(s, note)
    return s


s = new_slide()
oval_outline(s, SW - Inches(3.3), -Inches(1.5), Inches(4.6), Inches(4.6), LINE, 1.25)
oval_outline(s, SW - Inches(2.3), -Inches(0.5), Inches(2.6), Inches(2.6), C("1B2A47"), 1)
rect(s, SW - Inches(1.16), Inches(0.66), Inches(0.1), Inches(0.1), GREEN, radius=0.5)
tf = box(s, SW - Inches(2.6), Inches(0.58), Inches(1.35), Inches(0.26))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
run(p, "SYSTEM ONLINE", 10, GREEN, True)
tf = box(s, MARGIN, Inches(2.0), Inches(11), Inches(0.35))
run(tf.paragraphs[0], "HACKATHON PROTOTYPE · AUGUST 2026", 13, ACCENT, True)
tf = box(s, MARGIN, Inches(2.38), Inches(12), Inches(1.1))
run(tf.paragraphs[0], "PersonnelAI", 56, WHITE, True)
tf = box(s, MARGIN, Inches(3.55), Inches(9.6), Inches(0.9))
run(tf.paragraphs[0], "AI-Based Predictive Personnel Stress & Welfare Monitoring System",
    19, MUTED)
rect(s, MARGIN, Inches(4.42), Inches(1.1), Pt(3), ACCENT)
tf = box(s, MARGIN, Inches(4.68), Inches(11), Inches(0.5))
p = tf.paragraphs[0]
run(p, "\u201cFrom detecting risk\u2026 to helping people take the next step.\u201d", 16, TEXT, False, FONT, True)
chips_row(s, ["React 18", "FastAPI", "RandomForest", "Privacy-by-Design"],
          SH - Inches(1.05))
footer(s, 1)
notes(s, 'Hook: "We built an AI system that notices when a person is struggling '
         '\u2014 before they burn out \u2014 and automatically offers them supportive next '
         'steps, confidentially." Visual: landing page hero (radar + glass card).')

slide_content(
    2, "The problem", "Burnout is invisible until it breaks something",
    [
        ("Detected too late \u2014", "after performance drops or people quit."),
        ("Periodic, manual, subjective checks", "\u2014 easy to fake (\u201cI\u2019m fine\u201d)."),
        ("Stigma", "stops honest self-reporting, especially in uniformed / shift-based workforces."),
        ("No early-warning signal", "and no structured way for welfare teams to offer help."),
        ("Result:", "preventable attrition, sick leave, low morale, degraded readiness."),
    ],
    "Emphasize the gap: organizations collect almost no continuous, low-friction "
    "wellbeing signal \u2014 and even when they do, there is no safe process to act on it.",
    callout_text="Organizations react months late because nothing watches between annual surveys.",
    callout_color=ROSE,
)

slide_content(
    3, "Our solution", "A welfare intelligence platform, end to end",
    [
        ("Daily micro check-in \u2014", "under 2 minutes, 8 simple indicators."),
        ("AI risk prediction \u2014", "RandomForest scores every check-in instantly."),
        ("Explainable results \u2014", "every score shows which factors drove it, and why."),
        ("Automated support plan \u2014", "tiered, practical, voluntary actions from the top factors."),
        ("Officer console \u2014", "anonymized early warnings so humans stay in charge."),
        ("Confidential by design \u2014", "individuals are never exposed; only aggregates travel upward."),
    ],
    '"This isn\u2019t surveillance. It\u2019s an early-warning radar plus a support dispatcher '
    '\u2014 with the person, not the manager, at the center."',
    chips=["Check-in", "AI Prediction", "Contributing Factors", "Support Plan"],
    chip_colors=[ACCENT, C("60A5FA"), AMBER, GREEN],
)

slide_two_col(
    4, "Feature overview", "What ships in the prototype",
    "For the individual",
    [
        "Daily check-in \u2014 8 indicators + free-text note",
        "Instant risk indicator with confidence %",
        "Ranked contributing factors + narrative",
        "Welfare support plan with tiered actions",
        "Mark as Done \u00b7 Remind Me Later \u00b7 View Support",
        "Follow-up cadence + Schedule Reminder",
        "Wellness library \u00b7 trends dashboard",
        "Consent toggles \u00b7 full data export",
    ],
    "For the organization",
    [
        "Early Warning Center with severity triage",
        "Alert review \u2192 intervention workflow",
        "Interventions tracked to completion",
        "Anonymized unit analytics + reports",
        "Broadcast notifications",
        "Audit log with search (admin)",
        "Role-based access control",
        "Scenario demo mode \u2014 zero persistence",
    ],
    "Demo accounts: WELFARE-01 and ADMIN-01, password demo1234.",
)

slide_diagram(
    5, "System architecture", "Clean layers, swappable model",
    [
        "React 18 + Vite SPA   (TypeScript \u00b7 Tailwind \u00b7 framer-motion \u00b7 Recharts)",
        "\u2193",
        "FastAPI   \u2014   routers: auth \u00b7 personnel \u00b7 ai \u00b7 admin   \u00b7   JWT \u00b7 Pydantic v2",
        "\u2193",
        "Service layer   \u2014   ML engine wrapper \u00b7 support-plan rule engine \u00b7 notifications",
        "\u2193",
        "SQLAlchemy 2 ORM  \u2192  SQLite        RandomForest engine (joblib, loaded at startup)",
    ],
    [
        ("Swappable model \u2014", "the ML engine hides behind the service layer."),
        ("Typed contracts \u2014", "REST JSON mirrors Pydantic schemas \u2194 TypeScript interfaces."),
        ("Self-documenting \u2014", "Swagger docs auto-generated at /docs."),
    ],
    "Point out clean separation: frontend never talks to the model directly; "
    "the model can be replaced without touching either side.",
)

slide_table(
    6, "Tech stack", "Chosen for speed, safety, and clarity",
    ["Layer", "Technology", "Why"],
    [
        ["Frontend", "React 18 \u00b7 TypeScript \u00b7 Vite", "Fast dev loop, type-safe contracts"],
        ["Styling", "Tailwind CSS v3 (CSS-variable tokens)", "Dual-theme glassmorphism system"],
        ["Animation", "framer-motion v13", "Production motion, reduced-motion aware"],
        ["Charts", "Recharts", "Declarative, animated analytics"],
        ["Backend", "FastAPI \u00b7 Pydantic v2 \u00b7 Uvicorn", "Async, auto-documented REST"],
        ["ORM / DB", "SQLAlchemy 2.0 (typed) + SQLite", "Zero-config prototyping persistence"],
        ["ML", "scikit-learn RandomForest + joblib", "Robust tabular baseline"],
        ["Quality", "38-check smoke test suite", "Every endpoint verified after changes"],
    ],
    [Inches(1.7), Inches(5.3), Inches(5.3)],
    "If asked about scale: SQLite swaps to Postgres via a single SQLAlchemy URL change.",
    row_h=Inches(0.48),
    font_size=12,
)

slide_content(
    7, "User journey", "Two minutes a day, seven steps, zero friction",
    [
        ("Register / Login \u2014", "role-based accounts, validated signup."),
        ("Daily check-in \u2014", "mood, sleep, fatigue, workload, satisfaction, duty hours, overtime, breaks."),
        ("Instant prediction \u2014", "gauge shows risk level, score %, model confidence."),
        ("\u201cWhy this result?\u201d \u2014", "top factors with animated impact bars."),
        ("Support plan appears automatically \u2014", "Mark as Done \u00b7 Remind Me Later \u00b7 View Support."),
        ("Follow-up \u2014", "\u201cNext check-in recommended: 7 days\u201d + one-click reminder."),
        ("Recommendations history \u2014", "filterable by status, snooze-aware."),
    ],
    "Stress friction removal: the entire loop takes less than two minutes per day.",
    body_size=15,
)

slide_content(
    8, "The AI model", "RandomForest \u2014 robust, probabilistic, explainable",
    [
        ("Training data:", "24,000 synthetically generated records \u2014 prototype-safe, no real personal data."),
        ("Inputs (9 features):", "workload, fatigue, sleep quality, duty hours, overtime frequency, "
         "job satisfaction, rest-break quality, self-reported stress, recent workload change."),
        ("Risk score:", "risk_score = P(Moderate) + P(High) from class probabilities."),
        ("Thresholds:", "High \u2265 0.70 score (or p(High) \u2265 0.40) \u00b7 Moderate \u2265 0.34 \u00b7 Low otherwise."),
        ("Confidence:", "0.60 + 0.37 \u00d7 max class probability."),
        ("Validation:", "Macro-F1 0.816 \u00b1 0.005 across 5-fold cross-validation."),
        ("Also outputs:", "ranked factor sensitivities + templated narrative + disclaimer."),
    ],
    "Why RandomForest: strong tabular baseline, handles non-linear interactions, "
    "and gives probabilities we turn into calibrated-feeling scores.",
    chips=["24k records", "9 features", "F1 0.816", "<50 ms inference"],
    chip_colors=[ACCENT, ACCENT, GREEN, AMBER],
    body_size=14.5,
)

slide_content(
    9, "Explainable AI", "Every number on screen has a visible reason",
    [
        ("Per-factor sensitivity analysis \u2014", "each feature probed individually against a neutral baseline;"),
        ("", ""),
        ("Impact normalized to %", "with direction tags:"),
        ("", ""),
        ("Narrative auto-generated \u2014", "\u201cWorkload is currently elevating the welfare-risk indicator.\u201d"),
        ("UI surfaces \u2014", "top-3 factors on the dashboard, full ranked list on the prediction page."),
    ],
    "Judges care about black-box concerns \u2014 this is our answer. Every number on "
    "screen has a visible reason behind it.",
    body_size=14.5,
)

slide_table(
    10, "Support plan engine", "(risk level \u00d7 top factor) \u2192 concrete supportive actions",
    ["Condition", "Headline step", "Priority"],
    [
        ["Risk High + top factor Workload family", "\u201cReview current workload\u201d", ("HIGH PRIORITY", "rose")],
        ["Risk Moderate + top factor Fatigue", "\u201cEncourage adequate rest period\u201d", ("RECOMMENDED", "amber")],
        ["Sleep quality low", "Recovery time + sleep/wellness resources", ("REC / OPT", "amber")],
        ["Job satisfaction low", "Confidential consultation + feedback channels", ("REC / OPT", "amber")],
        ["Risk Low", "Continue regular wellbeing check-ins", ("OPTIONAL", "green")],
    ],
    [Inches(4.6), Inches(5.1), Inches(2.6)],
    "Factor families: Recent Workload Change / Overtime Frequency map to Workload; "
    "Duty Hours / Rest Break Quality map to Fatigue. Every step ships with sub-action "
    "checklist, timeframe, resource content, and tone guardrails: voluntary, "
    "confidential, never disciplinary. Follow-up: 7-day cadence with idempotent reminders.",
    row_h=Inches(0.56),
    font_size=12.5,
)

slide_content(
    11, "Demo mode", "The 60-second wow path for judges",
    [
        ("Three scenario cards \u2014", "Low / Moderate / High presets mapped exactly to the model contract."),
        ("Animated pipeline \u2014", "~2.5 s visualization of five analysis steps ending in a completion burst."),
        ("Result reveal \u2014", "zone-colored gauge arc, count-up score, confidence bar, shimmering factor bars."),
        ("Tiered action cards \u2014", "the same support plan logic that real users receive."),
        ("Zero persistence \u2014", "/ai/demo-predict runs the real model writing nothing to the database;"),
        ("", ""),
    ],
    "This is the 60-second wow path: click a card, watch the pipeline animate, "
    "land on a fully explained result.",
    callout_text="Verified: latest-prediction timestamp unchanged after demo runs \u2014 judge clicks cannot pollute user data.",
    callout_color=GREEN,
    body_size=14.5,
)

slide_two_col(
    12, "Officer & admin console", "The AI flags. A human decides. Everything is auditable.",
    "Welfare officers see",
    [
        "Early Warning Center \u2014 severity-counted alerts",
        "Review workflow: Confirm / No Action / Follow-up",
        "Confirming auto-creates an Intervention,",
        "assigns an officer, sets a due date, and",
        "notifies the member \u2014 fully voluntary",
        "Interventions pipeline to completion",
        "Aggregated, anonymized reports",
    ],
    "Administrators additionally get",
    [
        "User management (create / deactivate / reset)",
        "Broadcast notifications to everyone",
        "Audit log with search \u2014 every sensitive",
        "action recorded (who, what, when)",
        "",
        "Human-in-the-loop is mandatory by design:",
        "no automatic action ever targets a person",
    ],
    '"Human-in-the-loop is mandatory: the AI flags, a human decides, and every '
    'decision is auditable."',
)

slide_content(
    13, "Privacy & ethics", "We designed the system we\u2019d want used on us",
    [
        ("Voluntary \u2014", "participation optional; any recommendation can be dismissed or snoozed."),
        ("Consent controls \u2014", "check-ins / optional feedback / notifications toggles."),
        ("Confidentiality \u2014", "individual responses never shown to peers or managers; aggregates only."),
        ("Human review required \u2014", "output is a supportive indicator, never automatic action against a person."),
        ("Supportive tone \u2014", "\u201coffer\u201d, \u201cencourage\u201d, \u201coptional\u201d \u2014 never disciplinary framing."),
        ("Data ownership \u2014", "one-click full personal-data export."),
        ("Disclaimers everywhere \u2014", "\u201cnot a medical diagnosis.\u201d"),
        ("Synthetic-only data", "in the prototype \u2014 nothing real leaves the room."),
    ],
    'Say it plainly: "We designed the system we\u2019d want used on us." This slide wins '
    "ethics-heavy judging rubrics.",
    glyph="\u2713",
    glyph_color=GREEN,
    body_size=14,
)

slide_content(
    14, "Design system", "Glassmorphism craft, engineered accessibility",
    [
        ("Glassmorphism UI \u2014", "frosted panels, layered shadows, airy light theme + deep-navy command theme."),
        ("Token architecture \u2014", "one set of CSS variables drives both themes; components written once."),
        ("Motion language \u2014", "staggered reveals, count-ups, blur-to-sharp entrances, sliding nav pill, route crossfades."),
        ("Iconography \u2014", "consistent Lucide set throughout."),
        ("Accessibility \u2014", "respects OS reduced-motion; never color-only signaling."),
    ],
    "Toggle themes live during the demo \u2014 same components, zero rewrites.",
    chips=["Light + Dark", "Reduced-motion safe", "Mobile responsive", "Lucide icons"],
    chip_colors=[ACCENT, GREEN, AMBER, C("60A5FA")],
)

slide_content(
    15, "Engineering quality", "Prototype speed without prototype shortcuts",
    [
        ("Typed end-to-end \u2014", "Pydantic schemas \u2194 TypeScript interfaces mirror each other."),
        ("38-check smoke suite \u2014", "auth \u2192 check-in \u2192 prediction \u2192 recommendations \u2192 alerts \u2192 interventions \u2192 admin; self-resetting."),
        ("Safe schema evolution \u2014", "startup migrations add columns without data loss."),
        ("Rule-engine verification \u2014", "all four canonical scenarios validated."),
        ("Idempotent endpoints \u2014", "reminders, notification reads, alert reviews."),
    ],
    "Every feature was shipped with its safety net \u2014 the smoke suite ran after "
    "each change in this session.",
)

slide_content(
    16, "Impact & use cases", "Where this helps on day one",
    [
        ("Who benefits \u2014", "defense/uniformed services, healthcare, call centers, any shift-based workforce."),
        ("Early detection \u2014", "support arrives weeks earlier than annual surveys allow."),
        ("Normalized conversations \u2014", "daily check-ins make wellbeing routine, not exceptional."),
        ("Officer efficiency \u2014", "triage focuses on genuinely flagged cases, context attached."),
        ("Responsible-AI template \u2014", "explainability + human oversight + confidentiality working together."),
    ],
    "Pick the vertical matching your judges \u2014 defense, healthcare, or BPO/call centers.",
)

slide_table(
    17, "Roadmap", "From prototype to production",
    ["Horizon", "Item"],
    [
        [("NEXT", "sky"), "Time-series trend forecasting; unit-level heatmaps over time"],
        [("NEXT", "sky"), "Multilingual check-ins; progressive web app (mobile-first)"],
        [("LATER", "amber"), "Rostering / HR integrations feeding fatigue context"],
        [("LATER", "amber"), "Privacy-preserving aggregate benchmarking across units"],
        [("LATER", "amber"), "Model upgrades: gradient boosting, calibration, fairness audits"],
    ],
    [Inches(1.6), Inches(10.7)],
    "Roadmap order signals priorities: forecasting and mobile reach first, "
    "integrations second, model sophistication last.",
    row_h=Inches(0.54),
    font_size=13,
)

s = new_slide()
tfq = box(s, Inches(1.2), Inches(2.3), SW - Inches(2.4), Inches(1.9))
p = para(tfq, True)
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.25
run(p, "\u201cMost systems detect risk and stop there.", 26, WHITE, True, FONT, True)
p2 = tfq.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
p2.line_spacing = 1.25
r = p2.add_run()
r.text = "PersonnelAI detects risk "
rr = r.font
rr.size = Pt(26); rr.bold = True; rr.italic = True; rr.name = FONT; rr.color.rgb = WHITE
r2 = p2.add_run()
r2.text = "and helps the person take a supportive next step"
rf = r2.font
rf.size = Pt(26); rf.bold = True; rf.italic = True; rf.name = FONT; rf.color.rgb = ACCENT
r3 = p2.add_run()
r3.text = " \u2014 explained, voluntary, and confidential.\u201d"
rf3 = r3.font
rf3.size = Pt(26); rf3.bold = True; rf3.italic = True; rf3.name = FONT; rf3.color.rgb = WHITE
rect(s, SW / 2 - Inches(0.55), Inches(4.45), Inches(1.1), Pt(3), ACCENT)
tf = box(s, MARGIN, Inches(4.8), SW - MARGIN * 2, Inches(0.4))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "LIVE DEMO \u2014 Login \u2192 Dashboard \u2192 Try Demo \u2192 all three scenarios (\u224890 seconds)", 15, MUTED, True)
chips_row(s, ["WELFARE-01 / demo1234", "ADMIN-01 / demo1234"], SH - Inches(1.5),
          colors=[GREEN, AMBER], centered=True)
footer(s, 18)
notes(s, "Land the closing line slowly, then switch to the browser and run the "
         "demo path end to end.")

slide_table(
    19, "Appendix", "Q&A quick-reference facts",
    ["Question", "Answer"],
    [
        ["Model accuracy?", "Macro-F1 = 0.816 \u00b1 0.005 (5-fold CV on 24k synthetic records)"],
        ["Does demo pollute data?", "No \u2014 /ai/demo-predict persists nothing; verified by timestamps"],
        ["What if user skips days?", "Follow-up anchors to last check-in; nudge fires around day 7"],
        ["Can officers see individuals?", "Only via alert-review with explicit confirmation; dashboards aggregate"],
        ["Is it diagnosis?", "No \u2014 explicitly disclaimed; supportive indicator only"],
        ["Scale?", "SQLite prototype; swaps to Postgres via SQLAlchemy URL change"],
        ["Snooze behavior?", "\u201cRemind Me Later\u201d snoozes 2 days, then re-surfaces"],
        ["Numbers to memorize", "24k records \u00b7 9 features \u00b7 F1 0.816 \u00b7 <2 min check-in \u00b7 7-day cadence"],
    ],
    [Inches(3.3), Inches(9.0)],
    "Keep this slide hidden during the pitch; jump here only when judges ask.",
    row_h=Inches(0.47),
    font_size=11.5,
)

prs.core_properties.title = "PersonnelAI \u2014 Pitch Deck"
prs.core_properties.author = "Team PersonnelAI"
prs.save(OUT)
print("Saved:", OUT)

check = Presentation(str(OUT))
print("Slides:", len(check.slides.__iter__.__self__._sldIdLst) if False else len(check.slides._sldIdLst))
for i, sl in enumerate(check.slides, 1):
    texts = [sh.text_frame.paragraphs[0].runs[0].text
             for sh in sl.shapes if sh.has_text_frame and sh.text_frame.paragraphs[0].runs]
    print(i, texts[1] if len(texts) > 1 else (texts[0] if texts else "(visual only)"))
